import streamlit as st
import fitz
import os
import re
import faiss
import numpy as np
import requests
from datetime import datetime

from sentence_transformers import SentenceTransformer
from docx import Document
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.util import Pt

# -----------------------------------
# PAGE CONFIG
# -----------------------------------

st.set_page_config(
    page_title="AI PM Assistant",
    layout="wide"
)

# -----------------------------------
# CLEAN ENTERPRISE UI
# -----------------------------------

st.markdown(
    """
    <style>

    .stApp {
        background-color: #F8FAFC;
        color: #1E293B;
        font-family: 'Segoe UI', sans-serif;
    }

    .block-container {
        padding-top: 1.5rem;
        padding-bottom: 2rem;
        padding-left: 2rem;
        padding-right: 2rem;
    }

    h1 {
        color: #0F172A;
        font-size: 2.8rem;
        font-weight: 700;
    }

    h2, h3 {
        color: #1E293B;
        font-weight: 600;
    }

    div[data-testid="metric-container"] {

        background: white;

        border: 1px solid #E2E8F0;

        padding: 1rem;

        border-radius: 12px;
    }

    .stButton > button {

        background-color: #2563EB;

        color: white;

        border: none;

        border-radius: 8px;

        padding: 0.6rem 1rem;

        font-weight: 600;
    }

    .stButton > button:hover {

        background-color: #1D4ED8;

        color: white;
    }

    .stDownloadButton > button {

        background-color: #059669;

        color: white;

        border: none;

        border-radius: 8px;

        padding: 0.6rem 1rem;

        font-weight: 600;
    }

    .stTextArea textarea {

        background-color: white;

        color: #111827;

        border: 1px solid #CBD5E1;

        border-radius: 10px;
    }

    section[data-testid="stFileUploader"] {

        background-color: white;

        border: 1px dashed #CBD5E1;

        border-radius: 10px;

        padding: 1rem;
    }

    button[data-baseweb="tab"] {

        color: #334155 !important;

        font-weight: 600;

        font-size: 15px;
    }

    button[data-baseweb="tab"][aria-selected="true"] {

        color: #2563EB !important;

        border-bottom: 2px solid #2563EB;
    }

    .streamlit-expanderHeader {

        color: #1E293B;

        font-weight: 600;
    }

    </style>
    """,
    unsafe_allow_html=True
)

# -----------------------------------
# CREATE FOLDERS
# -----------------------------------

os.makedirs("repository", exist_ok=True)
os.makedirs("generated_fsds", exist_ok=True)
os.makedirs("generation_history", exist_ok=True)

# -----------------------------------
# EMBEDDING MODEL
# -----------------------------------

embedding_model = SentenceTransformer(
    "all-MiniLM-L6-v2"
)

# -----------------------------------
# PROMPT TEMPLATE HELPERS
# -----------------------------------

import json

os.makedirs("prompt_templates/fsd", exist_ok=True)
os.makedirs("user_templates/fsd", exist_ok=True)
os.makedirs("user_settings", exist_ok=True)

DEFAULT_FSD_TEMPLATES = {
    "Standard": "prompt_templates/fsd/standard.txt",
    "Detailed": "prompt_templates/fsd/detailed.txt",
    "Executive": "prompt_templates/fsd/executive.txt"
}

SETTINGS_FILE = "user_settings/default_templates.json"


def load_prompt_template(file_path):

    try:

        with open(file_path, "r") as f:
            return f.read()

    except Exception as e:

        st.error(f"Error loading template: {e}")

        return ""


def save_user_template(template_name, content):

    file_path = f"user_templates/fsd/{template_name}.txt"

    with open(file_path, "w") as f:
        f.write(content)


def get_user_templates():

    templates = []

    folder = "user_templates/fsd"

    for file in os.listdir(folder):

        if file.endswith(".txt"):

            templates.append(
                file.replace(".txt", "")
            )

    return templates


def load_default_template():

    if os.path.exists(SETTINGS_FILE):

        with open(SETTINGS_FILE, "r") as f:

            settings = json.load(f)

            return settings.get("fsd_default")

    return "Standard"


def save_default_template(template_name):

    settings = {
        "fsd_default": template_name
    }

    with open(SETTINGS_FILE, "w") as f:

        json.dump(settings, f)


# -----------------------------------
# SESSION STATE
# -----------------------------------

if "generated_fsd" not in st.session_state:
    st.session_state.generated_fsd = ""

if "approved_fsd" not in st.session_state:
    st.session_state.approved_fsd = ""

if "qa_generated" not in st.session_state:
    st.session_state.qa_generated = False

# -----------------------------------
# TITLE
# -----------------------------------

st.title("AI Product Copilot")

st.markdown(
    """
### AI-powered workflow for product documentation, QA, release readiness, and GTM content generation.

Requirement → AI Draft → PM Review → QA Test Cases → Release Notes → Support Notes → GTM Content → GTM PPT
"""
)
repository_files = os.listdir(
    "repository"
)

# -----------------------------------
# METRICS
# -----------------------------------

generated_files = os.listdir(
    "generated_fsds"
)

fsd_count = len([
    f for f in generated_files
    if f.startswith("FSD_")
])

qa_count = len([
    f for f in generated_files
    if f.startswith("QA_Test_Cases_")
])

release_count = len([
    f for f in generated_files
    if f.startswith("Release_Notes_")
])

support_count = len([
    f for f in generated_files
    if f.startswith("Support_Notes_")
])

gtm_count = len([
    f for f in generated_files
    if f.startswith("GTM_Content_")
])

ppt_count = len([
    f for f in generated_files
    if f.startswith("GTM_Presentation_")
])

col1, col2, col3 = st.columns(3)

with col1:
    st.metric(
        "Repository Files",
        len(repository_files)
    )

with col2:
    st.metric(
        "FSDs",
        fsd_count
    )

with col3:
    st.metric(
        "QA Docs",
        qa_count
    )

col4, col5, col6 = st.columns(3)

with col4:
    st.metric(
        "Release Notes",
        release_count
    )

with col5:
    st.metric(
        "Support Notes",
        support_count
    )

with col6:
    st.metric(
        "PPTs",
        ppt_count
    )
# -----------------------------------
# TABS
# -----------------------------------

tab1, tab2, tab3, tab4, tab5, tab6 = st.tabs([
    "FSD Generation",
    "QA Test Cases",
    "Release Notes",
    "Support Notes",
    "GTM Content",
    "PPT Generator"
])
# ===================================
# TAB 1
# ===================================

# ===================================
# TAB 1
# ===================================

with tab1:

    with st.expander("View Repository Files"):

        if repository_files:

            for file in repository_files:
                st.write(file)

        else:
            st.warning("Repository is empty")

    st.subheader("Upload Additional Documents")

    uploaded_files = st.file_uploader(
        "Upload PDFs or DOCX",
        type=["pdf", "docx"],
        accept_multiple_files=True
    )

    st.subheader("Enter Requirement")

    requirement = st.text_area(
        "Requirement",
        height=200
    )

    # -----------------------------------
    # TEMPLATE SECTION
    # -----------------------------------

    st.subheader("FSD Template")

    default_templates = list(
        DEFAULT_FSD_TEMPLATES.keys()
    )

    saved_templates = get_user_templates()

    template_options = (
        default_templates +
        saved_templates
    )

    default_template = load_default_template()

    default_index = 0

    if default_template in template_options:

        default_index = template_options.index(
            default_template
        )

    selected_template = st.selectbox(
        "Select Template",
        template_options,
        index=default_index
    )

    # -----------------------------------
    # LOAD TEMPLATE
    # -----------------------------------

    if selected_template in DEFAULT_FSD_TEMPLATES:

        template_content = load_prompt_template(
            DEFAULT_FSD_TEMPLATES[selected_template]
        )

    else:

        template_content = load_prompt_template(
            f"user_templates/fsd/{selected_template}.txt"
        )

    # -----------------------------------
    # EDITABLE INSTRUCTIONS
    # -----------------------------------

    editable_prompt = st.text_area(
        "AI Instructions",
        value=template_content,
        height=300
    )

    # -----------------------------------
    # SAVE TEMPLATE
    # -----------------------------------

    new_template_name = st.text_input(
        "Save As Template Name"
    )

    if st.button("Save Template"):

        if new_template_name:

            save_user_template(
                new_template_name,
                editable_prompt
            )

            st.success(
                "Template Saved Successfully"
            )

    # -----------------------------------
    # DEFAULT TEMPLATE
    # -----------------------------------

    if st.button("Set Selected Template As Default"):

        save_default_template(
            selected_template
        )

        st.success(
            f"{selected_template} set as default"
        )

    # -----------------------------------
    # GENERATE FSD
    # -----------------------------------

    analyze = st.button("Generate FSD")

    if analyze:

        if not requirement:

            st.error("Please enter requirement")

        else:

            # -----------------------------------
            # SAVE UPLOADED FILES
            # -----------------------------------

            if uploaded_files:

                for uploaded_file in uploaded_files:

                    file_path = os.path.join(
                        "repository",
                        uploaded_file.name
                    )

                    with open(file_path, "wb") as f:

                        f.write(
                            uploaded_file.getbuffer()
                        )

            repository_files = os.listdir(
                "repository"
            )

            all_chunks = []

            chunk_size = 800
            chunk_overlap = 100

            # -----------------------------------
            # READ REPOSITORY FILES
            # -----------------------------------

            for file in repository_files:

                file_path = os.path.join(
                    "repository",
                    file
                )

                text = ""

                try:

                    if file.endswith(".pdf"):

                        doc = fitz.open(file_path)

                        for page in doc:
                            text += page.get_text()

                    elif file.endswith(".docx"):

                        doc = Document(file_path)

                        for para in doc.paragraphs:
                            text += para.text + " "

                    clean_text = re.sub(
                        r'\s+',
                        ' ',
                        text
                    )

                    start = 0

                    while start < len(clean_text):

                        end = start + chunk_size

                        chunk = clean_text[start:end]

                        all_chunks.append({

                            "source": file,

                            "text": chunk
                        })

                        start += (
                            chunk_size - chunk_overlap
                        )

                except Exception as e:

                    st.error(
                        f"Error processing {file}: {e}"
                    )

            # -----------------------------------
            # CREATE EMBEDDINGS
            # -----------------------------------

            chunk_texts = [

                chunk["text"]

                for chunk in all_chunks
            ]

            embeddings = embedding_model.encode(
                chunk_texts
            )

            embeddings = np.array(
                embeddings
            ).astype("float32")

            dimension = embeddings.shape[1]

            index = faiss.IndexFlatL2(
                dimension
            )

            index.add(embeddings)

            # -----------------------------------
            # SEARCH RELEVANT CONTEXT
            # -----------------------------------

            query_embedding = embedding_model.encode(
                [requirement]
            )

            query_embedding = np.array(
                query_embedding
            ).astype("float32")

            distances, indices = index.search(
                query_embedding,
                5
            )

            final_context = ""

            reference_files = set()

            st.subheader(
                "Relevant Repository Matches"
            )

            for idx in indices[0]:

                match = all_chunks[idx]

                reference_files.add(
                    match["source"]
                )

                final_context += (
                    f"\nSource: {match['source']}\n"
                )

                final_context += (
                    match["text"] + "\n"
                )

                st.markdown("---")

                st.write(match["text"])

            final_context = final_context[:12000]

            # -----------------------------------
            # FINAL AI PROMPT
            # -----------------------------------

            prompt = f"""
Enterprise Repository Context:
{final_context}

Requirement:
{requirement}

Instructions:
{editable_prompt}
"""

            # -----------------------------------
            # GEMINI API
            # -----------------------------------

            api_key = st.secrets["GEMINI_API_KEY"]

            url = (
                "https://generativelanguage.googleapis.com/"
                f"v1/models/gemini-2.5-flash:generateContent?key={api_key}"
            )

            payload = {
                "contents": [
                    {
                        "parts": [
                            {
                                "text": prompt
                            }
                        ]
                    }
                ]
            }

            response = requests.post(
                url,
                json=payload
            )

            response_json = response.json()

            try:

                ai_output = response_json[
                    "candidates"
                ][0]["content"]["parts"][0]["text"]

                st.session_state.generated_fsd = ai_output

                # -----------------------------------
                # SAVE GENERATION HISTORY
                # -----------------------------------
                
                timestamp = datetime.now().strftime(
                    "%Y%m%d_%H%M%S"
                )
                existing_history = os.listdir(
                    "generation_history"
                )
                
                version_number = len(existing_history) + 1
                
                history_data = {

                    "version": version_number,
                    
                    "timestamp": timestamp,
                
                    "requirement": requirement,
                
                    "selected_template": selected_template,
                
                    "instructions": editable_prompt,
                
                    "generated_output": ai_output,
                
                }
                
                history_file = (
                    f"generation_history/FSD_{timestamp}.json"
                )
                
                with open(history_file, "w") as f:
                
                    json.dump(
                        history_data,
                        f,
                        indent=4
                    )

                st.subheader("Generated FSD")

                st.write(ai_output)

                # -----------------------------------
                # AI QUALITY EVALUATION
                # -----------------------------------

                evaluation_prompt = f"""
                You are an AI Product Documentation Reviewer.
                
                Evaluate the following FSD.
                
                FSD:
                {ai_output}
                
                Score each category out of 10.
                
                Return ONLY in this format:
                
                Completeness: X/10
                Business Clarity: X/10
                Validation Coverage: X/10
                Technical Depth: X/10
                Risk Coverage: X/10
                Overall Score: X/100
                
                Also provide:
                1 short strength
                1 short improvement suggestion
                """

                evaluation_payload = {
                    "contents": [
                        {
                            "parts": [
                                {
                                    "text": evaluation_prompt
                                }
                            ]
                        }
                    ]
                }

                evaluation_response = requests.post(
                    url,
                    json=evaluation_payload
                )

                evaluation_json = evaluation_response.json()

                try:

                    evaluation_output = evaluation_json[
                        "candidates"
                    ][0]["content"]["parts"][0]["text"]

                    st.subheader(
                        "AI Quality Evaluation"
                    )

                    st.info(
                        evaluation_output
                    )
                                        # -----------------------------------
                    # UPDATE HISTORY WITH EVALUATION
                    # -----------------------------------

                    history_data["evaluation"] = (
                        evaluation_output
                    )

                    with open(history_file, "w") as f:

                        json.dump(
                            history_data,
                            f,
                            indent=4
                        )
                except:

                    st.warning(
                        "Evaluation generation failed"
                    )

                st.subheader(
                    "Reference Documents Used"
                )

                for ref in reference_files:
                    st.write(f"- {ref}")

            except:

                st.error(
                    "AI generation failed"
                )

                st.write(response_json)
        # -----------------------------------
        # GENERATION HISTORY VIEWER
        # -----------------------------------

        st.subheader("Generation History")

        history_files = []

        if os.path.exists("generation_history"):

            history_files = sorted(
                os.listdir("generation_history"),
                reverse=True
            )

        if history_files:

            for history in history_files[:5]:

                history_path = os.path.join(
                    "generation_history",
                    history
                )

                try:

                    with open(history_path, "r") as f:

                        history_data = json.load(f)

                    version = history_data.get(
                        "version",
                        "Old"
                    )
                    
                    score_display = "No Score"
                    
                    if "evaluation" in history_data:
                    
                        evaluation_text = history_data["evaluation"]
                    
                        match = re.search(
                            r"Overall Score:\s*(.*)",
                            evaluation_text
                        )
                    
                        if match:
                    
                            score_display = match.group(1)
                    
                    with st.expander(
                        f"Version {version} | {history_data['selected_template']} | Score: {score_display}"
                    ):
                        st.write(
                            f"Requirement: {history_data['requirement']}"
                        )

                        st.write(
                            f"Template: {history_data['selected_template']}"
                        )

                        st.text_area(
                            "Instructions",
                            history_data["instructions"],
                            height=150,
                            disabled=True,
                            key=f"instr_{history}"
                        )

                        st.text_area(
                            "Generated Output",
                            history_data["generated_output"],
                            height=300,
                            disabled=True,
                            key=f"output_{history}"
                        )
                        
                        if "evaluation" in history_data:
                        
                            st.text_area(
                                "AI Evaluation",
                                history_data["evaluation"],
                                height=200,
                                disabled=True,
                                key=f"eval_{history}"
                            )

                except Exception as e:

                    st.error(
                        f"Error loading history: {e}"
                    )

        else:

            st.info(
                "No generation history available"
            )

        st.subheader("PM Review & Edit")

        edited_fsd = st.text_area(
            "Review / Edit Generated FSD",
            value=st.session_state.generated_fsd,
            height=500
        )

        if st.button("Approve FSD"):

            st.session_state.approved_fsd = edited_fsd

            st.success(
                "FSD Approved Successfully"
            )

        if st.session_state.approved_fsd:

            approved_fsd = st.session_state.approved_fsd

            if st.button("Export Approved FSD"):

                doc = Document()

                doc.add_heading(
                    "Functional Specification Document",
                    level=1
                )

                current_date = datetime.now().strftime(
                    "%d-%m-%Y %H:%M"
                )

                doc.add_paragraph(
                    f"Generated On: {current_date}"
                )

                doc.add_paragraph(
                    "Generated By: AI PM Assistant"
                )

                doc.add_paragraph(
                    approved_fsd
                )

                timestamp = datetime.now().strftime(
                    "%Y%m%d_%H%M%S"
                )

                file_name = (
                    f"generated_fsds/FSD_{timestamp}.docx"
                )

                doc.save(file_name)

                st.success(
                    "FSD Exported Successfully"
                )

                with open(file_name, "rb") as file:

                    st.download_button(
                        label="Download FSD",
                        data=file,
                        file_name=f"FSD_{timestamp}.docx",
                        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                    )

# ===================================
# TAB 2
# ===================================

with tab2:

    st.subheader(
        "QA Test Case Generation"
    )

    if st.session_state.approved_fsd:

        if st.button("Generate QA Test Cases"):

            approved_fsd = st.session_state.approved_fsd

            test_case_prompt = f"""
You are a QA Test Engineer.

Approved FSD:
{approved_fsd}

Generate:

1. Functional Test Cases
2. Validation Test Cases
3. Negative Test Cases
4. Regression Test Areas
5. Edge Cases
"""

            api_key = st.secrets["GEMINI_API_KEY"]

            url = (
                "https://generativelanguage.googleapis.com/"
                f"v1/models/gemini-2.5-flash:generateContent?key={api_key}"
            )

            payload = {
                "contents": [
                    {
                        "parts": [
                            {
                                "text": test_case_prompt
                            }
                        ]
                    }
                ]
            }

            response = requests.post(
                url,
                json=payload
            )

            response_json = response.json()

            try:

                test_cases = response_json[
                    "candidates"
                ][0]["content"]["parts"][0]["text"]

                st.session_state.qa_generated = True

                st.success(
                    "QA Test Cases Generated Successfully"
                )

                st.write(test_cases)

                qa_doc = Document()

                qa_doc.add_heading(
                    "QA Test Cases",
                    level=1
                )

                current_date = datetime.now().strftime(
                    "%d-%m-%Y %H:%M"
                )

                qa_doc.add_paragraph(
                    f"Generated On: {current_date}"
                )

                qa_doc.add_paragraph(
                    "Generated By: AI PM Assistant"
                )

                qa_doc.add_paragraph(
                    test_cases
                )

                qa_file_name = (
                    f"generated_fsds/QA_Test_Cases_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
                )

                qa_doc.save(qa_file_name)

                with open(qa_file_name, "rb") as file:

                    st.download_button(

                        label="Download QA Test Cases",

                        data=file,

                        file_name=os.path.basename(
                            qa_file_name
                        ),

                        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                    )

            except:

                st.error(
                    "QA generation failed"
                )

    else:

        st.warning(
            "Approve FSD before generating QA Test Cases"
        )

# ===================================
# TAB 3
# ===================================

with tab3:

    st.subheader("Release Notes")

    if st.session_state.qa_generated:

        if st.button(
            "Generate Release Notes"
        ):

            release_prompt = f"""
Generate enterprise release notes from this FSD.

FSD:
{st.session_state.approved_fsd}

Return:

1. Feature Summary
2. Key Changes
3. Business Impact
4. Important Notes
5. User Actions Required
"""

            api_key = st.secrets["GEMINI_API_KEY"]

            url = (
                "https://generativelanguage.googleapis.com/"
                f"v1/models/gemini-2.5-flash:generateContent?key={api_key}"
            )

            payload = {
                "contents": [
                    {
                        "parts": [
                            {
                                "text": release_prompt
                            }
                        ]
                    }
                ]
            }

            response = requests.post(
                url,
                json=payload
            )

            response_json = response.json()

            try:

                release_notes = response_json[
                    "candidates"
                ][0]["content"]["parts"][0]["text"]

                st.write(release_notes)

                release_doc = Document()

                release_doc.add_heading(
                    "Release Notes",
                    level=1
                )

                current_date = datetime.now().strftime(
                    "%d-%m-%Y %H:%M"
                )

                release_doc.add_paragraph(
                    f"Generated On: {current_date}"
                )

                release_doc.add_paragraph(
                    "Generated By: AI PM Assistant"
                )

                release_doc.add_paragraph(
                    release_notes
                )

                release_file_name = (
                    f"generated_fsds/Release_Notes_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
                )

                release_doc.save(release_file_name)

                with open(release_file_name, "rb") as file:

                    st.download_button(

                        label="Download Release Notes",

                        data=file,

                        file_name=os.path.basename(
                            release_file_name
                        ),

                        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                    )

            except:

                st.error(
                    "Release note generation failed"
                )

    else:

        st.warning(
            "Generate QA Test Cases before Release Notes"
        )
# ===================================
# TAB 4 - SUPPORT NOTES
# ===================================

with tab4:

    st.subheader("Support Notes")

    if st.session_state.qa_generated:

        if st.button(
            "Generate Support Notes"
        ):

            support_prompt = f"""
Generate enterprise support notes from this FSD.

FSD:
{st.session_state.approved_fsd}

Return:

1. Feature Overview
2. Configuration Dependencies
3. Known Limitations
4. Troubleshooting Notes
5. Common User Issues
6. Important Support Considerations
7. Impacted Areas
8. FAQs
"""

            api_key = st.secrets["GEMINI_API_KEY"]

            url = (
                "https://generativelanguage.googleapis.com/"
                f"v1/models/gemini-2.5-flash:generateContent?key={api_key}"
            )

            payload = {
                "contents": [
                    {
                        "parts": [
                            {
                                "text": support_prompt
                            }
                        ]
                    }
                ]
            }

            response = requests.post(
                url,
                json=payload
            )

            response_json = response.json()

            try:

                support_notes = response_json[
                    "candidates"
                ][0]["content"]["parts"][0]["text"]

                st.write(support_notes)

                # -----------------------------------
                # EXPORT SUPPORT NOTES
                # -----------------------------------

                support_doc = Document()

                support_doc.add_heading(
                    "Support Notes",
                    level=1
                )

                current_date = datetime.now().strftime(
                    "%d-%m-%Y %H:%M"
                )

                support_doc.add_paragraph(
                    f"Generated On: {current_date}"
                )

                support_doc.add_paragraph(
                    "Generated By: AI PM Assistant"
                )

                support_doc.add_paragraph(
                    support_notes
                )

                support_file_name = (
                    f"generated_fsds/Support_Notes_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
                )

                support_doc.save(support_file_name)

                with open(support_file_name, "rb") as file:

                    st.download_button(

                        label="Download Support Notes",

                        data=file,

                        file_name=os.path.basename(
                            support_file_name
                        ),

                        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                    )

            except:

                st.error(
                    "Support note generation failed"
                )

    else:

        st.warning(
            "Generate QA Test Cases before Support Notes"
        )
# ===================================
# TAB 5 - GTM CONTENT
# ===================================

with tab5:

    st.subheader("GTM Content Generation")

    if st.session_state.qa_generated:

        if st.button(
            "Generate GTM Content"
        ):

            gtm_prompt = f"""
Generate GTM content from this FSD.

FSD:
{st.session_state.approved_fsd}

Return:

1. Feature Highlights
2. Business Benefits
3. Customer Value Proposition
4. Sales Talking Points
5. Marketing Summary
6. Brochure Content
7. PPT Summary
8. Customer Announcement Draft
"""

            api_key = st.secrets["GEMINI_API_KEY"]

            url = (
                "https://generativelanguage.googleapis.com/"
                f"v1/models/gemini-2.5-flash:generateContent?key={api_key}"
            )

            payload = {
                "contents": [
                    {
                        "parts": [
                            {
                                "text": gtm_prompt
                            }
                        ]
                    }
                ]
            }

            response = requests.post(
                url,
                json=payload
            )

            response_json = response.json()

            try:

                gtm_content = response_json[
                    "candidates"
                ][0]["content"]["parts"][0]["text"]

                st.write(gtm_content)

                # -----------------------------------
                # EXPORT GTM CONTENT
                # -----------------------------------

                gtm_doc = Document()

                gtm_doc.add_heading(
                    "GTM Content",
                    level=1
                )

                current_date = datetime.now().strftime(
                    "%d-%m-%Y %H:%M"
                )

                gtm_doc.add_paragraph(
                    f"Generated On: {current_date}"
                )

                gtm_doc.add_paragraph(
                    "Generated By: AI PM Assistant"
                )

                gtm_doc.add_paragraph(
                    gtm_content
                )

                gtm_file_name = (
                    f"generated_fsds/GTM_Content_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
                )

                gtm_doc.save(gtm_file_name)

                with open(gtm_file_name, "rb") as file:

                    st.download_button(

                        label="Download GTM Content",

                        data=file,

                        file_name=os.path.basename(
                            gtm_file_name
                        ),

                        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                    )

            except:

                st.error(
                    "GTM content generation failed"
                )

    else:

        st.warning(
            "Generate QA Test Cases before GTM Content"
        )

# ===================================
# TAB 6 - PPT GENERATOR
# ===================================

with tab6:

    st.subheader("GTM PPT Generator")

    if st.session_state.qa_generated:

        if st.button(
            "Generate PPT"
        ):

            ppt_prompt = f"""
Generate presentation slide content from this FSD.

FSD:
{st.session_state.approved_fsd}

Return slide-wise content.

Format:

Slide 1: Title
Slide 2: Business Problem
Slide 3: Feature Overview
Slide 4: Key Functionalities
Slide 5: Business Benefits
Slide 6: Workflow Summary
Slide 7: Customer Value
Slide 8: Conclusion
"""

            api_key = st.secrets["GEMINI_API_KEY"]

            url = (
                "https://generativelanguage.googleapis.com/"
                f"v1/models/gemini-2.5-flash:generateContent?key={api_key}"
            )

            payload = {
                "contents": [
                    {
                        "parts": [
                            {
                                "text": ppt_prompt
                            }
                        ]
                    }
                ]
            }

            response = requests.post(
                url,
                json=payload
            )

            response_json = response.json()

            try:

                ppt_content = response_json[
                    "candidates"
                ][0]["content"]["parts"][0]["text"]

                st.write(ppt_content)

                # -----------------------------------
                # CREATE STYLED PPT
                # -----------------------------------

                prs = Presentation()

                slides_data = ppt_content.split(
                    "Slide "
                )

                for slide_data in slides_data:

                    if slide_data.strip() == "":
                        continue

                    lines = slide_data.split("\n")

                    title = lines[0].strip()

                    content = "\n".join(
                        lines[1:]
                    ).strip()

                    # -----------------------------------
                    # TITLE SLIDE
                    # -----------------------------------

                    if "1:" in title:

                        slide_layout = prs.slide_layouts[0]

                        slide = prs.slides.add_slide(
                            slide_layout
                        )

                        slide.shapes.title.text = (
                            "AI-Powered GTM Presentation"
                        )

                        subtitle = slide.placeholders[1]

                        subtitle.text = (
                            "Generated by AI PM Assistant"
                        )

                        title_shape = slide.shapes.title

                        title_para = (
                            title_shape.text_frame.paragraphs[0]
                        )

                        title_para.font.size = Pt(28)

                        title_para.font.bold = True

                        title_para.font.color.rgb = RGBColor(
                            37,
                            99,
                            235
                        )

                    # -----------------------------------
                    # CONTENT SLIDES
                    # -----------------------------------

                    else:

                        slide_layout = prs.slide_layouts[1]

                        slide = prs.slides.add_slide(
                            slide_layout
                        )

                        # TITLE

                        title_shape = slide.shapes.title

                        title_shape.text = title

                        title_para = (
                            title_shape.text_frame.paragraphs[0]
                        )

                        title_para.font.size = Pt(24)

                        title_para.font.bold = True

                        title_para.font.color.rgb = RGBColor(
                            15,
                            23,
                            42
                        )

                        # CONTENT

                        body_shape = slide.placeholders[1]

                        text_frame = body_shape.text_frame

                        text_frame.text = ""

                        for line in content.split("\n"):

                            if line.strip() == "":
                                continue

                            p = text_frame.add_paragraph()

                            p.text = line.strip()

                            p.font.size = Pt(16)

                            p.font.color.rgb = RGBColor(
                                51,
                                65,
                                85
                            )

                # -----------------------------------
                # SAVE PPT
                # -----------------------------------

                ppt_file_name = (
                    f"generated_fsds/GTM_Presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
                )

                prs.save(ppt_file_name)

                st.success(
                    "PPT Generated Successfully"
                )

                with open(
                    ppt_file_name,
                    "rb"
                ) as file:

                    st.download_button(

                        label="Download PPT",

                        data=file,

                        file_name=os.path.basename(
                            ppt_file_name
                        ),

                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )

            except Exception as e:

                st.error(
                    f"PPT generation failed: {e}"
                )

    else:

        st.warning(
            "Generate QA Test Cases before PPT Generation"
        )
